from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION

from pathlib import Path

PALETTE = {
    "bg": RGBColor(245, 247, 250),
    "primary": RGBColor(33, 150, 243),
    "secondary": RGBColor(99, 199, 126),
    "accent": RGBColor(255, 193, 7),
    "ink": RGBColor(30, 41, 59),
    "muted": RGBColor(100, 116, 139),
    "white": RGBColor(255, 255, 255),
    "card": RGBColor(255, 255, 255),
}

FONTS = {
    "title": ("Montserrat", Pt(44)),
    "subtitle": ("Inter", Pt(20)),
    "h2": ("Montserrat", Pt(28)),
    "body": ("Inter", Pt(18)),
    "caption": ("Inter", Pt(12)),
}

MARGINS = {
    "left": Inches(0.7),
    "right": Inches(0.7),
    "top": Inches(0.6),
    "bottom": Inches(0.6),
}

def set_slide_background(slide, color=PALETTE["bg"]):
    left = top = Inches(0)
    width = Inches(13.333)
    height = Inches(7.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()

def add_header_footer(slide, title_text=None, page_number=None, logo_path=None):
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(0.5)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PALETTE["white"]
    bar.line.fill.background()

    if title_text:
        tx = slide.shapes.add_textbox(MARGINS["left"], Inches(0.08), Inches(8), Inches(0.34))
        p = tx.text_frame.paragraphs[0]
        p.text = title_text
        p.font.name, p.font.size = FONTS["caption"]
        p.font.color.rgb = PALETTE["muted"]
        p.alignment = PP_ALIGN.LEFT

    if page_number is not None:
        tx = slide.shapes.add_textbox(Inches(12.6), Inches(0.08), Inches(0.6), Inches(0.34))
        p = tx.text_frame.paragraphs[0]
        p.text = str(page_number)
        p.font.name, p.font.size = FONTS["caption"]
        p.font.color.rgb = PALETTE["muted"]
        p.alignment = PP_ALIGN.RIGHT

    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(logo_path, Inches(12.7), Inches(0.06), height=Inches(0.38))

def add_title_slide(prs, title, subtitle="", author="", logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PALETTE["bg"])

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = PALETTE["primary"]
    band.line.fill.background()

    tx = slide.shapes.add_textbox(MARGINS["left"], Inches(2.3), Inches(9), Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name, p.font.size = FONTS["title"]
    p.font.bold = True
    p.font.color.rgb = PALETTE["ink"]

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.space_before = Pt(10)
        p2.font.name, p2.font.size = FONTS["subtitle"]
        p2.font.color.rgb = PALETTE["muted"]

    if author:
        tx2 = slide.shapes.add_textbox(MARGINS["left"], Inches(4.6), Inches(9), Inches(0.6))
        p3 = tx2.text_frame.paragraphs[0]
        p3.text = author
        p3.font.name, p3.font.size = FONTS["caption"]
        p3.font.color.rgb = PALETTE["muted"]

    add_header_footer(slide, title_text="")
    return slide

def add_section_slide(prs, heading, blurb=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PALETTE["bg"])
    add_header_footer(slide, title_text=heading, page_number=len(prs.slides))

    tx = slide.shapes.add_textbox(MARGINS["left"], Inches(1.7), Inches(10), Inches(1.6))
    p = tx.text_frame.paragraphs[0]
    p.text = heading
    p.font.name, p.font.size = FONTS["h2"]
    p.font.bold = True
    p.font.color.rgb = PALETTE["ink"]

    if blurb:
        p2 = tx.text_frame.add_paragraph()
        p2.text = blurb
        p2.font.name, p2.font.size = FONTS["body"]
        p2.font.color.rgb = PALETTE["muted"]
        p2.space_before = Pt(8)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGINS["left"], Inches(2.8), Inches(1.5), Inches(0.12))
    line.fill.solid()
    line.fill.fore_color.rgb = PALETTE["accent"]
    line.line.fill.background()

    return slide

def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PALETTE["bg"])
    add_header_footer(slide, title_text=title, page_number=len(prs.slides))

    tx = slide.shapes.add_textbox(MARGINS["left"], Inches(1.1), Inches(10), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name, p.font.size = FONTS["h2"]
    p.font.bold = True
    p.font.color.rgb = PALETTE["ink"]

    tx2 = slide.shapes.add_textbox(MARGINS["left"], Inches(2.0), Inches(11.4), Inches(4.8))
    tf = tx2.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
        para.level = 0
        para.font.name, para.font.size = FONTS["body"]
        para.font.color.rgb = PALETTE["ink"]
        para.space_after = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_points, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PALETTE["bg"])
    add_header_footer(slide, title_text=title, page_number=len(prs.slides))

    tx = slide.shapes.add_textbox(MARGINS["left"], Inches(1.1), Inches(10), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name, p.font.size = FONTS["h2"]
    p.font.bold = True
    p.font.color.rgb = PALETTE["ink"]

    txL = slide.shapes.add_textbox(MARGINS["left"], Inches(2.0), Inches(5.5), Inches(4.8))
    tfL = txL.text_frame
    tfL.word_wrap = True
    for i, b in enumerate(left_points):
        para = tfL.paragraphs[0] if i == 0 else tfL.add_paragraph()
        para.text = b
        para.level = 0
        para.font.name, para.font.size = FONTS["body"]
        para.font.color.rgb = PALETTE["ink"]
        para.space_after = Pt(6)

    txR = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.9))
    tfR = txR.text_frame
    tfR.word_wrap = True
    for i, b in enumerate(right_points):
        para = tfR.paragraphs[0] if i == 0 else tfR.add_paragraph()
        para.text = b
        para.level = 0
        para.font.name, para.font.size = FONTS["body"]
        para.font.color.rgb = RGBColor(30, 41, 59)
        para.space_after = Pt(6)

    return slide

def add_chart_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PALETTE["bg"])
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.chart import XL_LABEL_POSITION

    tx = slide.shapes.add_textbox(MARGINS["left"], Inches(1.1), Inches(10), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name, p.font.size = FONTS["h2"]
    p.font.bold = True
    p.font.color.rgb = PALETTE["ink"]

    chart_data = ChartData()
    chart_data.categories = ["A", "B", "C", "D"]
    chart_data.add_series("2025", (4.3, 2.7, 3.8, 1.9))

    x, y, cx, cy = Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.9)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart.chart_style = 2
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True

    for s in chart.series:
        s.data_labels.show_value = True
        s.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    return slide

def build_deck(slides_spec, path="output_beautiful.pptx", title="", subtitle="", author="", logo_path=None):
    prs = Presentation()
    if title:
        add_title_slide(prs, title, subtitle, author, logo_path)
    for spec in slides_spec:
        kind = spec.get("type", "bullets")
        if kind == "section":
            add_section_slide(prs, spec.get("heading", "Section"), spec.get("blurb", ""))
        elif kind == "two-column":
            add_two_column_slide(prs, spec.get("title", "Two column"), spec.get("left", []), spec.get("right", []))
        elif kind == "chart":
            add_chart_slide(prs, spec.get("title", "Chart"))
        else:
            add_bullets_slide(prs, spec.get("title", "Slide"), spec.get("bullets", []))
    prs.save(path)
    return path

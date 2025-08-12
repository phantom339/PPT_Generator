"""
ppt_generator.py
Beautiful, image-enabled PPT generation with Gemini/Gemma + Pexels + python-pptx.

Requires:
  pip install google-generativeai python-pptx pillow python-dotenv requests

.env:
  GOOGLE_API_KEY=<key>       # or GEMINI_API_KEY=<key>
  GEMINI_MODEL=gemma-3-27b-it  # or gemini-1.5-flash etc.
  PEXELS_API_KEY=<key>

Companion:
  beauty_ppt.py  (design helpers & theme constants)
"""

import os
import re
import io
import json
import shutil
import tempfile
from dataclasses import dataclass
from typing import List, Optional, Dict, Any

import requests
from PIL import Image
from dotenv import load_dotenv
import google.generativeai as genai

# Use design helpers + theme constants
from beauty_ppt import (
    PALETTE, FONTS, MARGINS,
    add_title_slide, add_section_slide, add_bullets_slide,
    add_two_column_slide, add_chart_slide, set_slide_background, add_header_footer
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# -------------------- ENV --------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY") or ""
MODEL_NAME = os.getenv("GEMINI_MODEL", "gemma-3-27b-it")
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "")

# -------------------- Types --------------------
@dataclass
class DeckRequest:
    topic: str
    audience: str = "General"
    tone: str = "Crisp, clear, professional"
    visual_style: str = "Modern, clean, minimal"
    slide_count: int = 8
    title: Optional[str] = None
    subtitle: Optional[str] = None
    author: Optional[str] = None
    logo_path: Optional[str] = None

# -------------------- Prompts --------------------
JSON_SCHEMA_HINT = r"""
Return ONLY valid JSON (no markdown fences). Schema:

{
  "slides": [
    // choose appropriate slide types and include image_query when helpful
    { "type": "section", "heading": "string", "blurb": "string (optional)" },
    { "type": "bullets", "title": "string", "bullets": ["string", "..."], "image_query": "string (optional)" },
    { "type": "two-column", "title": "string", "left": ["string", "..."], "right": ["string", "..."], "image_query": "string (optional)" },
    { "type": "image-hero", "caption": "string (optional)", "image_query": "string" }
  ],
  "meta": { "title": "string", "subtitle": "string (optional)" }
}

Rules:
- Aim for the requested slide count (±1 ok).
- Bullets short (<=50 words).
- Include image_query on ~40–60% slides where a visual helps.
- One opener (section), 3–5 bullets slides, 1 two-column, 1 summary.
"""

def build_gemini_prompt(req: DeckRequest) -> str:
    title = req.title or req.topic
    return f"""
You are a senior presentation strategist.
Create a concise deck outline with tasteful visuals.

Topic: {req.topic}
Audience: {req.audience}
Tone: {req.tone}
Visual: {req.visual_style}
Desired Slides: {req.slide_count}

Include helpful image ideas as 'image_query' text (e.g., 'data center racks', 'AI workflow diagram', 'wind turbines at sunset').

Main deck title: "{title}"
Subtitle: "{req.subtitle or ''}"

{JSON_SCHEMA_HINT}
""".strip()

# -------------------- AI + JSON utils --------------------
def configure_genai():
    if not GOOGLE_API_KEY:
        raise ValueError("Missing GOOGLE_API_KEY / GEMINI_API_KEY in .env")
    genai.configure(api_key=GOOGLE_API_KEY)
    return genai.GenerativeModel(MODEL_NAME)

def safe_json_parse(text: str) -> Dict[str, Any]:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"```(json)?", "", text, flags=re.IGNORECASE).strip()
        text = re.sub(r"```$", "", text).strip()
    m = re.search(r"\{.*\}", text, flags=re.DOTALL)
    if m:
        text = m.group(0)
    return json.loads(text)

def generate_outline(req: DeckRequest) -> Dict[str, Any]:
    model = configure_genai()
    resp = model.generate_content(build_gemini_prompt(req))
    return safe_json_parse((resp.text or "").strip())

# -------------------- Pexels image search --------------------
PEXELS_SEARCH_URL = "https://api.pexels.com/v1/search"

def fetch_image_bytes(query: str, orientation: str = "landscape") -> Optional[bytes]:
    if not PEXELS_API_KEY:
        return None
    headers = {"Authorization": PEXELS_API_KEY}
    params = {"query": query, "per_page": 1, "orientation": orientation}
    try:
        r = requests.get(PEXELS_SEARCH_URL, headers=headers, params=params, timeout=15)
        r.raise_for_status()
        data = r.json()
        photos = data.get("photos", [])
        if not photos:
            return None
        url = photos[0]["src"]["large"] or photos[0]["src"].get("original")
        img = requests.get(url, timeout=15)
        img.raise_for_status()
        return img.content
    except Exception:
        return None

def save_image_tmp(img_bytes: bytes, tmpdir: str, fname_prefix="img") -> Optional[str]:
    try:
        path = os.path.join(tmpdir, f"{fname_prefix}.jpg")
        with Image.open(io.BytesIO(img_bytes)) as im:
            # Ensure RGB JPEG
            rgb = im.convert("RGB")
            rgb.save(path, format="JPEG", quality=88, optimize=True)
        return path
    except Exception:
        return None

# -------------------- Image slide helpers (styled) --------------------
def add_image_right_slide(prs, title: str, bullets: List[str], image_path: str):
    """Two-column layout with text left and image right, matching theme."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PALETTE["bg"])
    add_header_footer(slide, title_text=title, page_number=len(prs.slides))

    # Title
    tx = slide.shapes.add_textbox(MARGINS["left"], Inches(1.0), Inches(10), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name, p.font.size = FONTS["h2"]
    p.font.bold = True
    p.font.color.rgb = PALETTE["ink"]

    # Bullets (left)
    tx2 = slide.shapes.add_textbox(MARGINS["left"], Inches(1.8), Inches(5.6), Inches(4.9))
    tf = tx2.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets or []):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
        para.level = 0
        para.font.name, para.font.size = FONTS["body"]
        para.font.color.rgb = PALETTE["ink"]

    # Image (right)
    try:
        slide.shapes.add_picture(image_path, Inches(7.2), Inches(1.8), height=Inches(4.9))
    except Exception:
        pass
    return slide

def add_image_hero_slide(prs, image_path: str, caption: str = ""):
    """Full-bleed hero image with subtle caption bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, PALETTE["bg"])
    add_header_footer(slide, title_text="", page_number=len(prs.slides))

    # Full image area
    try:
        slide.shapes.add_picture(image_path, Inches(0.0), Inches(0.5), width=Inches(13.333))
    except Exception:
        # fallback: colored block if image fails
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.5), Inches(13.333), Inches(6.5))
        rect.fill.solid(); rect.fill.fore_color.rgb = PALETTE["secondary"]; rect.line.fill.background()

    # Caption strip
    if caption:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(6.6), Inches(13.333), Inches(0.8))
        bar.fill.solid(); bar.fill.fore_color.rgb = PALETTE["white"]; bar.line.fill.background()
        tx = slide.shapes.add_textbox(MARGINS["left"], Inches(6.65), Inches(10), Inches(0.7))
        p = tx.text_frame.paragraphs[0]
        p.text = caption
        p.font.name, p.font.size = FONTS["body"]
        p.font.color.rgb = PALETTE["muted"]
    return slide

# -------------------- Core generator --------------------
class PPTGenerator:
    def __init__(self, api_key: Optional[str] = None, model: Optional[str] = None):
        self.api_key = (api_key or GOOGLE_API_KEY).strip()
        if not self.api_key:
            raise ValueError("GOOGLE_API_KEY / GEMINI_API_KEY is missing.")
        self.model_name = (model or MODEL_NAME).strip()

    def _outline(self, req: DeckRequest) -> Dict[str, Any]:
        return generate_outline(req)

    def generate_presentation(
        self,
        topic: str,
        audience: str = "General",
        tone: str = "Crisp, clear, professional",
        visual_style: str = "Modern, clean, minimal",
        slide_count: int = 8,
        title: Optional[str] = None,
        subtitle: Optional[str] = None,
        author: Optional[str] = None,
        logo_path: Optional[str] = None,
        output_path: str = "output_beautiful.pptx",
        download_images: bool = True,
    ) -> str:
        req = DeckRequest(
            topic=topic, audience=audience, tone=tone, visual_style=visual_style,
            slide_count=slide_count, title=title or topic, subtitle=subtitle,
            author=author, logo_path=logo_path
        )

        tmpdir = tempfile.mkdtemp(prefix="ppt_imgs_")
        prs = Presentation()
        try:
            try:
                data = self._outline(req)
                meta = data.get("meta", {}) or {}
                deck_title = meta.get("title") or req.title or topic
                deck_subtitle = meta.get("subtitle") or req.subtitle or ""
                slides = data.get("slides", [])
            except Exception as e:
                # Minimal fallback
                deck_title, deck_subtitle = (req.title or topic, req.subtitle or "")
                slides = [
                    {"type": "section", "heading": topic, "blurb": audience},
                    {"type": "bullets", "title": "Overview", "bullets": ["What it is", "Why it matters", "Key points"]},
                    {"type": "two-column", "title": "Pros & Cons",
                     "left": ["Pro 1", "Pro 2"], "right": ["Con 1", "Con 2"]},
                    {"type": "bullets", "title": "Next Steps", "bullets": ["Do X", "Measure Y", "Review Z"]},
                ]

            # Title slide
            add_title_slide(prs, deck_title, deck_subtitle, author or "", logo_path)

            # Slide builder
            for s in slides:
                stype = s.get("type", "bullets")
                image_query = s.get("image_query")
                image_path = None

                # fetch image if requested/available
                if download_images and image_query:
                    img_bytes = fetch_image_bytes(image_query)
                    if img_bytes:
                        image_path = save_image_tmp(img_bytes, tmpdir, fname_prefix=f"img_{len(prs.slides)}")

                if stype == "section":
                    add_section_slide(prs, s.get("heading", "Section"), s.get("blurb", ""))

                elif stype == "bullets":
                    title_txt = s.get("title", "Slide")
                    bullets = s.get("bullets", [])
                    if image_path:
                        add_image_right_slide(prs, title_txt, bullets, image_path)
                    else:
                        add_bullets_slide(prs, title_txt, bullets)

                elif stype == "two-column":
                    title_txt = s.get("title", "Two column")
                    left = s.get("left", [])
                    right = s.get("right", [])
                    # If we have an image, show left bullets + right image; else standard two-column text
                    if image_path:
                        add_image_right_slide(prs, title_txt, left or right, image_path)
                    else:
                        add_two_column_slide(prs, title_txt, left, right)

                elif stype == "image-hero":
                    cap = s.get("caption", "")
                    if image_path:
                        add_image_hero_slide(prs, image_path, cap)
                    else:
                        # If no image fetched, just make a section w/ caption
                        add_section_slide(prs, cap or "Visual", "")

                elif stype == "chart":
                    add_chart_slide(prs, s.get("title", "Chart"))

                else:
                    # Unknown => fallback to bullets
                    add_bullets_slide(prs, s.get("title", "Slide"), s.get("bullets", []))

            prs.save(output_path)
            return output_path

        finally:
            # cleanup
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except Exception:
                pass


# -------------------- CLI (optional) --------------------
if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Generate a beautiful, image-enabled PPT from a topic.")
    parser.add_argument("--topic", required=True)
    parser.add_argument("--audience", default="General")
    parser.add_argument("--tone", default="Crisp, clear, professional")
    parser.add_argument("--style", default="Modern, clean, minimal")
    parser.add_argument("--slides", type=int, default=8)
    parser.add_argument("--title", default=None)
    parser.add_argument("--subtitle", default=None)
    parser.add_argument("--author", default=None)
    parser.add_argument("--logo", default=None)
    parser.add_argument("--out", default="output_beautiful.pptx")
    parser.add_argument("--no-images", action="store_true")
    args = parser.parse_args()

    gen = PPTGenerator()
    out = gen.generate_presentation(
        topic=args.topic,
        audience=args.audience,
        tone=args.tone,
        visual_style=args.style,
        slide_count=args.slides,
        title=args.title,
        subtitle=args.subtitle,
        author=args.author,
        logo_path=args.logo,
        output_path=args.out,
        download_images=not args.no_images,
    )
    print(f"Saved: {out}")
